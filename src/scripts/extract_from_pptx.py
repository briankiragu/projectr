"""Script to get track lyrics as json data from .PPTX"""


import argparse
import json
import logging
import os
import re
import structlog
import sys
import unicodedata
import uuid

from pptx import Presentation
from typing import Set


# Create an instance of the logger.
structlog.configure(
    wrapper_class=structlog.make_filtering_bound_logger(logging.INFO),
    processors=[
        structlog.processors.add_log_level,
        structlog.processors.TimeStamper(fmt="iso"),
        structlog.dev.ConsoleRenderer(),
    ]
)
log = structlog.get_logger()


def replace_with_ascii(phrase: str) -> str:
    """Replace the tricky unicode characters"""

    return re.sub(
        r"[\u201c\u201d]",
        '"',
        phrase.replace("\u2019", "'").replace("\u2018", "").replace("\u2013", "-")
    )

def extract_title_from_string(title: str) -> str:
    """Remove formatting from titles"""

    return unicodedata.normalize(
        "NFKC", replace_with_ascii(re.sub(r"[\u000b\n]", " ", title))
    ).strip()


def extract_lyrics_from_string(verse: str) -> str:
    """From a single string, extract each lyric as an item in an array."""

    return unicodedata.normalize(
        "NFKC", replace_with_ascii(re.sub(r"[\u000b\n]", "\n", verse))
    ).strip()


def extract_data_from_pptx(input_pptx: str):
    log.info("Attempting to extract data from .pptx", file=input_pptx)
    prs = Presentation(input_pptx)
    data = []

    for slide in prs.slides:
        slide_data = [
            shape.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text.strip()
        ]
        data.append(slide_data)

    return data


def format_extracted_lyrics(
    tenant_id: str,
    slides: list[list[str]]
) -> list[dict[str, int | str | list[str]]]:
    """Format the extracted data to the expected format"""

    titles: list[str] = []
    lyrics: list[dict[str, int | str | list[str] | None]] = []
    current_title: str = ""
    current_content: list[str] = []

    for slide in slides:
        if len(slide) >= 2:
            if current_title not in titles:
                titles.append(current_title)
                lyrics.append(
                    {
                        "tenant": tenant_id,
                        "title": current_title,
                        "content": "\n\n".join(current_content),
                        "status": "published",
                        "sort": None
                    }
                )

            current_title = extract_title_from_string(slide[0])
            current_content = [
                extract_lyrics_from_string(*slide[1:])
            ]
        elif len(slide) == 1:
            current_content.append(
                extract_lyrics_from_string(*slide)
            )

    return lyrics  # type: ignore

try:
    # Configure the argument parse
    parser = argparse.ArgumentParser()
    parser.add_argument("-t", "--tenant", help="Tenant ID")
    parser.add_argument("-p", "--path", nargs='*', action='append',  help="Path to the .pptx file")
    args = parser.parse_args()

    if args.tenant == None:
        raise Exception("No tenant ID provided. Use the -t or --tenant option to provide a tenant ID.")

    if args.path == None:
        raise Exception("No paths to .pptx file(s). Use the -p or --path option to provide a relative path to a .pptx file.")

    # Find the directory in which the current script resides:
    dirname = os.path.dirname(os.path.realpath(__file__))

    # Get the Tenant ID from the CLI.
    tenant_id = args.tenant
    log.debug(f"Tenant ID provided.", tenant_id=tenant_id)

    # Get all the .PPTX files supplied in the CLI.
    file_paths = [path for row in args.path for path in row]
    log.debug(f"File paths provided.", file_paths=file_paths)

    # Create a variable to hold all the data.
    extracted_data = []
    log.info(f"Attempting to extract data...",)
    for file_path in file_paths:
        extracted_data += extract_data_from_pptx(file_path)
        log.debug(
            f"Extracted data from pptx file.",
            records=len(extracted_data)
        )

    # Format the data as JSON.
    formatted_json_data: list[
        dict[str, int | str | list[str]]
    ] = format_extracted_lyrics(tenant_id, extracted_data)

    # Set the name of the output file.
    output_json = f"{dirname}/data/json/content.formatted.json"
    log.info(f"Attempting to format data...")

    # Write the json output file.
    with open(output_json, "w", encoding="utf8") as json_file:
        json.dump(formatted_json_data, json_file)
    log.info("Lyrics data formatted successfully and saved.", output=output_json)
except Exception as e:
    log.error(e)
